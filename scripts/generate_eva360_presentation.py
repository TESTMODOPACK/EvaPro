#!/usr/bin/env python3
"""Generate the EVA360 institutional presentation (PowerPoint).

Brand identity: Ascenda — black + gold gradient. 16:9 widescreen.
Output: docs/presentaciones/EVA360-Presentacion-Institucional.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Brand palette (Ascenda / EVA360)
# ---------------------------------------------------------------------------
BG_BLACK = RGBColor(0x0B, 0x0B, 0x0F)        # Deep black background
BG_PANEL = RGBColor(0x16, 0x16, 0x1C)        # Slightly lighter panel
GOLD_PRIMARY = RGBColor(0xD4, 0xAF, 0x37)    # Ascenda gold
GOLD_LIGHT = RGBColor(0xF4, 0xD8, 0x7A)      # Light gold for highlights
GOLD_DARK = RGBColor(0x8C, 0x6E, 0x1F)       # Dark gold accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED = RGBColor(0xC8, 0xC8, 0xCF)
TEXT_DIM = RGBColor(0x8A, 0x8A, 0x92)
LINE_GOLD = RGBColor(0xB8, 0x95, 0x2A)

FONT_TITLE = "Calibri"   # Inter substitute available everywhere
FONT_BODY = "Calibri"

ROOT = Path(__file__).resolve().parent.parent
LOGO_PATH = ROOT / "docs" / "eva360_logo_principal_transparent.png"
OUT_DIR = ROOT / "docs" / "presentaciones"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "EVA360-Presentacion-Institucional.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=FONT_BODY,
    size=18,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines,
    *,
    font=FONT_BODY,
    size=14,
    color=TEXT_MUTED,
    bullet_color=GOLD_PRIMARY,
    line_spacing=1.25,
):
    """lines: list of strings. Lines starting with '• ' get a gold bullet."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if line.startswith("• "):
            r1 = p.add_run()
            r1.text = "● "
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = bullet_color
            r2 = p.add_run()
            r2.text = line[2:]
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_gold_bar(slide, left, top, width, height=Emu(45720)):  # 0.05" tall
    bar = add_rect(slide, left, top, width, height, GOLD_PRIMARY)
    return bar


def slide_header(slide, eyebrow: str, title: str, subtitle: str | None = None):
    """Standard top header on content slides."""
    add_text(
        slide,
        Inches(0.6),
        Inches(0.35),
        Inches(8),
        Inches(0.3),
        eyebrow.upper(),
        size=11,
        bold=True,
        color=GOLD_PRIMARY,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(0.65),
        Inches(12),
        Inches(0.7),
        title,
        font=FONT_TITLE,
        size=30,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.6),
            Inches(1.25),
            Inches(12),
            Inches(0.45),
            subtitle,
            size=14,
            color=TEXT_MUTED,
            italic=True,
        )
    add_gold_bar(slide, Inches(0.6), Inches(1.7), Inches(2.5))


def slide_footer(slide, page_num: int, total: int):
    add_rect(
        slide,
        Inches(0),
        Inches(7.15),
        Inches(13.333),
        Inches(0.05),
        GOLD_DARK,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(7.25),
        Inches(8),
        Inches(0.3),
        "EVA360 · Ascenda SpA · Confidencial",
        size=9,
        color=TEXT_DIM,
    )
    add_text(
        slide,
        Inches(11.7),
        Inches(7.25),
        Inches(1.2),
        Inches(0.3),
        f"{page_num:02d} / {total:02d}",
        size=9,
        bold=True,
        color=GOLD_PRIMARY,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, left, top, width, height, title, body_lines, *, title_size=14, body_size=11):
    add_rect(slide, left, top, width, height, BG_PANEL, line_color=GOLD_DARK)
    add_rect(slide, left, top, Inches(0.08), height, GOLD_PRIMARY)
    add_text(
        slide,
        left + Inches(0.25),
        top + Inches(0.15),
        width - Inches(0.4),
        Inches(0.4),
        title,
        size=title_size,
        bold=True,
        color=GOLD_LIGHT,
    )
    add_paragraphs(
        slide,
        left + Inches(0.25),
        top + Inches(0.6),
        width - Inches(0.4),
        height - Inches(0.7),
        body_lines,
        size=body_size,
        color=TEXT_MUTED,
        line_spacing=1.2,
    )


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = 14


# ---- Slide 1 — Cover ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)

# Decorative gold bars (left side, ascending)
bar_left = Inches(0.8)
heights = [0.5, 0.9, 1.4, 2.0, 2.7, 3.5, 4.4]
gap = 0.18
bar_w = 0.42
for i, h in enumerate(heights):
    add_rect(
        s,
        bar_left + Inches(i * (bar_w + gap)),
        Inches(6.0 - h),
        Inches(bar_w),
        Inches(h),
        GOLD_PRIMARY if i == len(heights) - 1 else GOLD_DARK,
    )

# Right-side text block
add_text(
    s,
    Inches(6.0),
    Inches(2.2),
    Inches(7),
    Inches(0.4),
    "ASCENDA SPA · 2026",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_text(
    s,
    Inches(6.0),
    Inches(2.6),
    Inches(7),
    Inches(1.4),
    "EVA360",
    font=FONT_TITLE,
    size=88,
    bold=True,
    color=GOLD_LIGHT,
)
add_gold_bar(s, Inches(6.05), Inches(4.05), Inches(3.5))
add_text(
    s,
    Inches(6.0),
    Inches(4.2),
    Inches(7),
    Inches(0.6),
    "Gestión del desempeño con IA",
    font=FONT_TITLE,
    size=24,
    bold=True,
    color=WHITE,
)
add_text(
    s,
    Inches(6.0),
    Inches(4.7),
    Inches(7),
    Inches(0.5),
    "Una sola plataforma. Todo el ciclo de personas.",
    size=16,
    color=TEXT_MUTED,
    italic=True,
)
add_text(
    s,
    Inches(6.0),
    Inches(6.1),
    Inches(7),
    Inches(0.4),
    "Presentación institucional · Nuevos usuarios + Inversionistas",
    size=11,
    color=TEXT_DIM,
)
add_text(
    s,
    Inches(6.0),
    Inches(6.5),
    Inches(7),
    Inches(0.4),
    "ricardo@ascenda.cl  ·  ascenda.cl  ·  Santiago, Chile",
    size=11,
    bold=True,
    color=GOLD_PRIMARY,
)


# ---- Slide 2 — El problema ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "01 · Contexto",
    "El problema que resolvemos",
    "Las empresas evalúan a sus personas mal, tarde y caro.",
)
add_paragraphs(
    s,
    Inches(0.6),
    Inches(2.1),
    Inches(7.5),
    Inches(4.5),
    [
        "• 78% de los empleados desconfía de la evaluación tradicional (Gallup, 2024).",
        "• Procesos manuales en Excel y Forms generan sesgo, baja participación y datos no accionables.",
        "• Soluciones enterprise (Workday, Lattice, Culture Amp) cuestan USD $20–50/empleado/mes — inalcanzables para pymes LATAM.",
        "• Resultado: rotación 25% mayor y planes de desarrollo (PDI) que nunca se ejecutan.",
        "• Líderes toman decisiones a ciegas; RR.HH. invierte semanas armando reportes que nadie lee.",
    ],
    size=15,
    line_spacing=1.5,
)
# Insight panel
card(
    s,
    Inches(8.5),
    Inches(2.1),
    Inches(4.3),
    Inches(4.5),
    "INSIGHT CLAVE",
    [
        "El problema no es la falta de herramientas.",
        "",
        "Es que las herramientas existentes no fueron diseñadas para empresas latinoamericanas de 50–500 empleados.",
        "",
        "EVA360 nació para ese vacío.",
    ],
    title_size=13,
    body_size=13,
)
slide_footer(s, 2, TOTAL_SLIDES)


# ---- Slide 3 — Qué es EVA360 ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "02 · Solución",
    "Qué es EVA360",
    "Una plataforma SaaS B2B integral. Reemplaza Excel + 4 herramientas separadas.",
)

modules = [
    ("Evaluaciones 360°", "Configurables por cargo, área y nivel."),
    ("OKRs y objetivos", "Con seguimiento, calibración y check-ins."),
    ("Feedback continuo", "1-on-1 estructurados + agendas IA."),
    ("PDI personalizados", "Generados con IA a partir del perfil."),
    ("Encuestas de clima", "Mood check-ins semanales de 10 segundos."),
    ("Reconocimientos", "Badges, kudos y reconocimiento social."),
    ("Insights con IA", "Detección de riesgo de fuga + brechas."),
    ("Calibración + 9-box", "Analytics de equipo y rotación."),
    ("Reclutamiento + ATS", "Postulantes integrados al ciclo."),
]

cols = 3
card_w = Inches(4.05)
card_h = Inches(1.35)
gap_x = Inches(0.15)
gap_y = Inches(0.18)
start_x = Inches(0.6)
start_y = Inches(2.05)
for idx, (mod_title, mod_body) in enumerate(modules):
    r = idx // cols
    c = idx % cols
    left = start_x + c * (card_w + gap_x)
    top = start_y + r * (card_h + gap_y)
    add_rect(s, left, top, card_w, card_h, BG_PANEL, line_color=GOLD_DARK)
    add_rect(s, left, top, card_w, Inches(0.06), GOLD_PRIMARY)
    add_text(
        s,
        left + Inches(0.2),
        top + Inches(0.18),
        card_w - Inches(0.3),
        Inches(0.4),
        mod_title,
        size=14,
        bold=True,
        color=GOLD_LIGHT,
    )
    add_text(
        s,
        left + Inches(0.2),
        top + Inches(0.62),
        card_w - Inches(0.3),
        Inches(0.7),
        mod_body,
        size=11,
        color=TEXT_MUTED,
    )
slide_footer(s, 3, TOTAL_SLIDES)


# ---- Slide 4 — Producto operativo hoy ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "03 · Producto",
    "Lo que ya está funcionando hoy",
    "EVA360 no es una promesa: es producto operativo en producción.",
)
metrics = [
    ("39", "Controllers backend", "NestJS 11"),
    ("30+", "Páginas UI activas", "Next.js 14"),
    ("78", "Entidades de dominio", "PostgreSQL + RLS"),
    ("100%", "Multi-tenant aislado", "Row-Level Security"),
]
mw = Inches(2.95)
mh = Inches(1.7)
for i, (num, label, sub) in enumerate(metrics):
    left = Inches(0.6) + i * (mw + Inches(0.1))
    top = Inches(2.05)
    add_rect(s, left, top, mw, mh, BG_PANEL, line_color=GOLD_DARK)
    add_text(
        s,
        left,
        top + Inches(0.2),
        mw,
        Inches(0.8),
        num,
        font=FONT_TITLE,
        size=46,
        bold=True,
        color=GOLD_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.05),
        mw,
        Inches(0.35),
        label,
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.35),
        mw,
        Inches(0.3),
        sub,
        size=10,
        color=TEXT_DIM,
        align=PP_ALIGN.CENTER,
        italic=True,
    )
add_text(
    s,
    Inches(0.6),
    Inches(4.0),
    Inches(12),
    Inches(0.4),
    "Capacidades en producción",
    size=14,
    bold=True,
    color=GOLD_LIGHT,
)
add_paragraphs(
    s,
    Inches(0.6),
    Inches(4.45),
    Inches(6.0),
    Inches(2.5),
    [
        "• Integración productiva con Anthropic Claude SDK",
        "• PWA instalable + Push notifications nativas",
        "• 2FA + JWT + auditoría cross-tenant",
        "• Suscripciones, planes y add-ons IA con Stripe + MercadoPago",
    ],
    size=12,
    line_spacing=1.4,
)
add_paragraphs(
    s,
    Inches(7.0),
    Inches(4.45),
    Inches(6.0),
    Inches(2.5),
    [
        "• Documentación legal completa (T&C, DPA, SLA, NDA, Privacidad)",
        "• Multi-tenant con RLS validado en PostgreSQL",
        "• Observabilidad con Sentry + métricas internas",
        "• Demo público: eva360.ascenda.cl/demo",
    ],
    size=12,
    line_spacing=1.4,
)
slide_footer(s, 4, TOTAL_SLIDES)


# ---- Slide 5 — Experiencia por persona ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "04 · Usuarios",
    "Diseñado para tres audiencias",
    "Una experiencia diferente por rol, sobre una sola plataforma.",
)
personas = [
    (
        "EL COLABORADOR",
        [
            "• Dashboard personal con evaluaciones, objetivos y PDI.",
            "• Mood check-ins de 10 segundos por semana.",
            "• Reconocimientos públicos para celebrar logros.",
            "• Notificaciones inteligentes — sin spam.",
        ],
    ),
    (
        "EL LÍDER",
        [
            "• Vista 360° del equipo: desempeño, riesgo, atrasos.",
            "• Streaks de liderazgo: hábitos sanos gamificados.",
            "• 1-on-1 con agenda automática + notas IA.",
            "• Calibración asistida con 9-box.",
        ],
    ),
    (
        "RECURSOS HUMANOS",
        [
            "• Ciclos lanzados en minutos, no semanas.",
            "• Plantillas configurables por cargo y nivel.",
            "• Insights agregados: brechas, DEI, rotación, NPS.",
            "• Auditoría completa de cada calificación.",
        ],
    ),
]
pw = Inches(4.0)
ph = Inches(4.5)
for i, (title, lines) in enumerate(personas):
    left = Inches(0.6) + i * (pw + Inches(0.15))
    top = Inches(2.05)
    add_rect(s, left, top, pw, ph, BG_PANEL, line_color=GOLD_DARK)
    add_rect(s, left, top, pw, Inches(0.5), GOLD_PRIMARY)
    add_text(
        s,
        left,
        top + Inches(0.08),
        pw,
        Inches(0.4),
        title,
        size=13,
        bold=True,
        color=BG_BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_paragraphs(
        s,
        left + Inches(0.25),
        top + Inches(0.75),
        pw - Inches(0.4),
        ph - Inches(0.9),
        lines,
        size=12,
        line_spacing=1.45,
    )
slide_footer(s, 5, TOTAL_SLIDES)


# ---- Slide 6 — IA con Claude ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "05 · Diferencial técnico",
    "IA con Claude — embebida, no pegada",
    "La IA está en los flujos críticos del negocio, no en un chatbot lateral.",
)
ia_uses = [
    ("Análisis de feedback", "Convierte texto libre en competencias accionables."),
    ("PDI personalizado", "Genera planes de desarrollo desde fortalezas y brechas."),
    ("Riesgo de fuga", "Alertas tempranas con señales de mood + desempeño."),
    ("Resúmenes ejecutivos", "Cierre de ciclo en una página, automático."),
    ("Sugerencias OKR", "Objetivos alineados con la estrategia organizacional."),
    ("Análisis DEI", "Detección de sesgo en evaluaciones y promociones."),
]
for i, (t, b) in enumerate(ia_uses):
    r, c = i // 2, i % 2
    left = Inches(0.6) + c * Inches(6.2)
    top = Inches(2.1) + r * Inches(1.35)
    add_rect(s, left, top, Inches(6.0), Inches(1.2), BG_PANEL, line_color=GOLD_DARK)
    add_rect(s, left, top, Inches(0.06), Inches(1.2), GOLD_PRIMARY)
    add_text(
        s,
        left + Inches(0.25),
        top + Inches(0.15),
        Inches(5.6),
        Inches(0.4),
        t,
        size=14,
        bold=True,
        color=GOLD_LIGHT,
    )
    add_text(
        s,
        left + Inches(0.25),
        top + Inches(0.55),
        Inches(5.6),
        Inches(0.6),
        b,
        size=11,
        color=TEXT_MUTED,
    )
add_text(
    s,
    Inches(0.6),
    Inches(6.5),
    Inches(12),
    Inches(0.4),
    "Modelo de costos: créditos IA por uso. El cliente paga solo lo que consume.",
    size=12,
    italic=True,
    color=GOLD_PRIMARY,
)
slide_footer(s, 6, TOTAL_SLIDES)


# ---- Slide 7 — Arquitectura y seguridad ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "06 · Arquitectura",
    "Stack moderno + seguridad enterprise",
    "Multi-tenant real con aislamiento a nivel base de datos.",
)
stack_rows = [
    ("Frontend", "Next.js 14 + React + TypeScript", "PWA, SEO, performance"),
    ("Backend", "NestJS 11 + TypeORM", "API REST documentada"),
    ("Base de datos", "PostgreSQL + Row-Level Security", "Aislamiento multi-tenant"),
    ("IA", "Anthropic Claude SDK", "No entrena con datos del cliente"),
    ("Pagos", "Stripe + MercadoPago", "Chile + LATAM listo"),
    ("Auth", "JWT + 2FA + SSO (roadmap)", "Compatible Azure AD / Okta"),
    ("Observabilidad", "Sentry + métricas internas", "SLA medible"),
    ("Cumplimiento", "DPA + Privacidad + RLS", "GDPR-ready"),
]
header_y = Inches(2.05)
row_h = Inches(0.5)
table_left = Inches(0.6)
col_w = [Inches(2.4), Inches(5.0), Inches(5.0)]

# Header
header_x = table_left
for i, h in enumerate(["CAPA", "STACK", "GARANTÍA"]):
    add_rect(s, header_x, header_y, col_w[i], row_h, GOLD_DARK)
    add_text(
        s,
        header_x + Inches(0.2),
        header_y + Inches(0.12),
        col_w[i] - Inches(0.3),
        row_h,
        h,
        size=12,
        bold=True,
        color=BG_BLACK,
    )
    header_x += col_w[i]

for r, (a, b, c) in enumerate(stack_rows):
    y = header_y + (r + 1) * row_h
    fill = BG_PANEL if r % 2 == 0 else BG_BLACK
    x = table_left
    for i, cell in enumerate([a, b, c]):
        add_rect(s, x, y, col_w[i], row_h, fill, line_color=GOLD_DARK)
        is_bold = i == 0
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(0.13),
            col_w[i] - Inches(0.3),
            row_h,
            cell,
            size=11,
            bold=is_bold,
            color=GOLD_LIGHT if is_bold else TEXT_MUTED,
        )
        x += col_w[i]
slide_footer(s, 7, TOTAL_SLIDES)


# ---- Slide 8 — Mercado ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "07 · Mercado",
    "Oportunidad: HR Tech LATAM",
    "Mercado creciendo a 17% CAGR, con clara brecha en mid-market.",
)
tam_blocks = [
    ("TAM", "USD $4.8B", "HR Tech LATAM (2025)\nGartner / Statista"),
    ("SAM", "USD $620M", "SaaS desempeño LATAM\nmid-market · Mordor"),
    ("SOM", "USD $18M", "Chile + México + Colombia\n3 años · bottom-up"),
]
for i, (label, num, desc) in enumerate(tam_blocks):
    left = Inches(0.6) + i * Inches(4.15)
    top = Inches(2.1)
    add_rect(s, left, top, Inches(3.95), Inches(2.4), BG_PANEL, line_color=GOLD_DARK)
    add_rect(s, left, top, Inches(3.95), Inches(0.45), GOLD_PRIMARY)
    add_text(
        s,
        left,
        top + Inches(0.05),
        Inches(3.95),
        Inches(0.4),
        label,
        size=14,
        bold=True,
        color=BG_BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(0.7),
        Inches(3.95),
        Inches(0.8),
        num,
        font=FONT_TITLE,
        size=38,
        bold=True,
        color=GOLD_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left + Inches(0.2),
        top + Inches(1.6),
        Inches(3.55),
        Inches(0.7),
        desc,
        size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
add_text(
    s,
    Inches(0.6),
    Inches(4.85),
    Inches(12),
    Inches(0.4),
    "Cliente ideal (ICP)",
    size=14,
    bold=True,
    color=GOLD_LIGHT,
)
add_paragraphs(
    s,
    Inches(0.6),
    Inches(5.25),
    Inches(12),
    Inches(2),
    [
        "• 50–500 empleados · sectores servicios, retail, fintech, salud privada.",
        "• Geografía: Chile, México, Colombia, Perú.",
        "• Pain validado: ya intentaron Excel/Forms y fracasaron.",
        "• Trigger comercial: nueva regulación laboral (Ley 40 horas, NCh 3262 en Chile).",
    ],
    size=12,
    line_spacing=1.4,
)
slide_footer(s, 8, TOTAL_SLIDES)


# ---- Slide 9 — Competencia ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "08 · Competencia",
    "Posicionamiento competitivo",
    "Precio LATAM, IA nativa, multi-tenant moderno.",
)
comp_rows = [
    ("Lattice (US)", "USD $11", "Limitada", "No", "No"),
    ("Culture Amp (AU)", "USD $14", "No", "No", "No"),
    ("Bamboo HR", "USD $9", "No", "No", "Parcial"),
    ("Buk (CL)", "UF 0.4 + extras", "No", "Sí", "Sí"),
    ("Rankmi (CL)", "A consultar", "Limitada", "Sí", "Sí"),
    ("EVA360", "CLP $1.500–4.000", "Claude nativa", "Sí (RLS)", "es / en / pt"),
]
header_y = Inches(2.0)
row_h = Inches(0.42)
table_left = Inches(0.6)
col_w = [Inches(2.8), Inches(2.5), Inches(2.2), Inches(2.5), Inches(2.1)]
hdrs = ["COMPETIDOR", "PRECIO/EMPL./MES", "IA", "MULTI-TENANT LATAM", "IDIOMA"]

x = table_left
for i, h in enumerate(hdrs):
    add_rect(s, x, header_y, col_w[i], row_h, GOLD_DARK)
    add_text(
        s,
        x + Inches(0.15),
        header_y + Inches(0.08),
        col_w[i],
        row_h,
        h,
        size=11,
        bold=True,
        color=BG_BLACK,
    )
    x += col_w[i]

for r, row in enumerate(comp_rows):
    y = header_y + (r + 1) * row_h
    is_eva = row[0] == "EVA360"
    fill = GOLD_DARK if is_eva else (BG_PANEL if r % 2 == 0 else BG_BLACK)
    x = table_left
    for i, cell in enumerate(row):
        add_rect(s, x, y, col_w[i], row_h, fill, line_color=GOLD_DARK)
        bold = is_eva or i == 0
        color = BG_BLACK if is_eva else (GOLD_LIGHT if i == 0 else TEXT_MUTED)
        add_text(
            s,
            x + Inches(0.15),
            y + Inches(0.08),
            col_w[i],
            row_h,
            cell,
            size=11,
            bold=bold,
            color=color,
        )
        x += col_w[i]

add_text(
    s,
    Inches(0.6),
    Inches(5.3),
    Inches(12),
    Inches(0.4),
    "Tres ventajas defendibles",
    size=14,
    bold=True,
    color=GOLD_LIGHT,
)
add_paragraphs(
    s,
    Inches(0.6),
    Inches(5.7),
    Inches(12),
    Inches(1.4),
    [
        "• Precio 3–5x menor que competidores US, ajustado al poder de compra LATAM.",
        "• IA generativa nativa (no add-on caro): insights automáticos en feedback, competencias y PDI.",
        "• Stack moderno + multi-tenant que escala sin rehacer la arquitectura.",
    ],
    size=12,
    line_spacing=1.4,
)
slide_footer(s, 9, TOTAL_SLIDES)


# ---- Slide 10 — Modelo de negocio ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "09 · Negocio",
    "Modelo de negocio",
    "SaaS B2B con planes mensuales/anuales por número de empleados.",
)
plans = [
    ("Starter", "Hasta 50", "$200.000", "$2.4M ARR"),
    ("Growth", "Hasta 200", "$800.000", "$9.6M ARR"),
    ("Business", "Hasta 500", "$2.000.000", "$24M ARR"),
    ("Enterprise", "500+", "Custom (UF)", "$50M+ ARR"),
]
pw = Inches(2.95)
ph = Inches(2.4)
for i, (name, emp, price, arr) in enumerate(plans):
    left = Inches(0.6) + i * (pw + Inches(0.1))
    top = Inches(2.0)
    is_highlighted = name == "Business"
    fill = GOLD_DARK if is_highlighted else BG_PANEL
    add_rect(s, left, top, pw, ph, fill, line_color=GOLD_PRIMARY if is_highlighted else GOLD_DARK)
    add_text(
        s,
        left,
        top + Inches(0.2),
        pw,
        Inches(0.4),
        name.upper(),
        size=13,
        bold=True,
        color=BG_BLACK if is_highlighted else GOLD_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(0.6),
        pw,
        Inches(0.35),
        emp,
        size=11,
        color=BG_BLACK if is_highlighted else TEXT_DIM,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.05),
        pw,
        Inches(0.6),
        price,
        font=FONT_TITLE,
        size=22,
        bold=True,
        color=BG_BLACK if is_highlighted else GOLD_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.5),
        pw,
        Inches(0.3),
        "CLP / mes",
        size=10,
        color=BG_BLACK if is_highlighted else TEXT_DIM,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.85),
        pw,
        Inches(0.3),
        arr,
        size=11,
        bold=True,
        color=BG_BLACK if is_highlighted else GOLD_PRIMARY,
        align=PP_ALIGN.CENTER,
    )

# Unit economics
add_text(
    s,
    Inches(0.6),
    Inches(4.7),
    Inches(12),
    Inches(0.4),
    "Unit economics objetivo · mes 18",
    size=14,
    bold=True,
    color=GOLD_LIGHT,
)
ue = [
    ("CAC", "$800.000"),
    ("LTV", "$14.000.000"),
    ("LTV/CAC", "17.5x"),
    ("Payback", "4 meses"),
    ("Gross margin", "78%"),
]
uw = Inches(2.4)
for i, (k, v) in enumerate(ue):
    left = Inches(0.6) + i * (uw + Inches(0.1))
    top = Inches(5.2)
    add_rect(s, left, top, uw, Inches(1.4), BG_PANEL, line_color=GOLD_DARK)
    add_text(
        s,
        left,
        top + Inches(0.2),
        uw,
        Inches(0.4),
        k,
        size=11,
        bold=True,
        color=GOLD_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(0.55),
        uw,
        Inches(0.7),
        v,
        font=FONT_TITLE,
        size=22,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
slide_footer(s, 10, TOTAL_SLIDES)


# ---- Slide 11 — Tracción ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "10 · Tracción",
    "Estado actual y próximos hitos",
    "Producto listo. Comercial en aceleración.",
)
add_text(
    s,
    Inches(0.6),
    Inches(2.05),
    Inches(6),
    Inches(0.4),
    "PRODUCTO",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_paragraphs(
    s,
    Inches(0.6),
    Inches(2.45),
    Inches(6.0),
    Inches(3.5),
    [
        "• Plataforma operativa: 39 controllers + 30 páginas UI.",
        "• 78 entidades de dominio modeladas.",
        "• Integración productiva con Anthropic Claude.",
        "• Documentación legal completa.",
        "• PWA + push + 2FA + auditoría.",
        "• Multi-tenant con RLS validado.",
    ],
    size=12,
    line_spacing=1.5,
)
add_text(
    s,
    Inches(7.0),
    Inches(2.05),
    Inches(6),
    Inches(0.4),
    "COMERCIAL · ACTUALIZAR AL PRESENTAR",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_paragraphs(
    s,
    Inches(7.0),
    Inches(2.45),
    Inches(6.0),
    Inches(3.5),
    [
        "• [N] tenants demo activos",
        "• [N] clientes pagados (pilots o suscripción)",
        "• [MRR] CLP MRR",
        "• [N] leads en pipeline desde landing pública",
        "• [N] integraciones con partners",
    ],
    size=12,
    line_spacing=1.5,
)
add_rect(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(1.3), BG_PANEL, line_color=GOLD_DARK)
add_rect(s, Inches(0.6), Inches(5.7), Inches(0.06), Inches(1.3), GOLD_PRIMARY)
add_text(
    s,
    Inches(0.85),
    Inches(5.85),
    Inches(11.5),
    Inches(0.4),
    "PRÓXIMOS 12 MESES",
    size=12,
    bold=True,
    color=GOLD_LIGHT,
)
add_text(
    s,
    Inches(0.85),
    Inches(6.25),
    Inches(11.5),
    Inches(0.7),
    "SSO SAML + SCIM  ·  Slack + Teams  ·  ISO 27001 (Etapa 1)  ·  Expansión piloto México",
    size=13,
    bold=True,
    color=WHITE,
)
slide_footer(s, 11, TOTAL_SLIDES)


# ---- Slide 12 — Roadmap ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "11 · Roadmap",
    "Roadmap a 12 meses",
    "8 hitos verificables · KPI de cierre: 15+ clientes, MRR USD $8.000+, NPS 80.",
)
roadmap_rows = [
    ("Q1", "SSO SAML + SCIM enterprise", "Compatible Azure AD, Okta, Google Workspace"),
    ("Q2", "Integraciones Slack + Teams", "App publicada en Slack Marketplace"),
    ("Q1–Q2", "Landing + SEO + content", "5 case studies + 1.000 sesiones/mes"),
    ("Q2–Q3", "Outbound + ABM mid-market", "10 demos calificadas/mes + 3 pilots"),
    ("Q3", "ISO 27001 (Etapa 1)", "Auditoría inicial completada"),
    ("Q3–Q4", "Expansión México", "2 clientes en CDMX cerrados"),
    ("Q4", "App nativa iOS / Android", "Publicación en App Store + Play Store"),
    ("Q4", "API pública + Marketplace", "3 integraciones de partners activas"),
]
header_y = Inches(2.0)
row_h = Inches(0.45)
table_left = Inches(0.6)
col_w = [Inches(1.5), Inches(4.5), Inches(6.6)]
hdrs = ["TRIMESTRE", "HITO", "RESULTADO VERIFICABLE"]

x = table_left
for i, h in enumerate(hdrs):
    add_rect(s, x, header_y, col_w[i], row_h, GOLD_DARK)
    add_text(
        s,
        x + Inches(0.15),
        header_y + Inches(0.1),
        col_w[i],
        row_h,
        h,
        size=11,
        bold=True,
        color=BG_BLACK,
    )
    x += col_w[i]

for r, (q, hito, res) in enumerate(roadmap_rows):
    y = header_y + (r + 1) * row_h
    fill = BG_PANEL if r % 2 == 0 else BG_BLACK
    x = table_left
    for i, cell in enumerate([q, hito, res]):
        add_rect(s, x, y, col_w[i], row_h, fill, line_color=GOLD_DARK)
        if i == 0:
            color = GOLD_PRIMARY
            bold = True
        elif i == 1:
            color = GOLD_LIGHT
            bold = True
        else:
            color = TEXT_MUTED
            bold = False
        add_text(
            s,
            x + Inches(0.15),
            y + Inches(0.1),
            col_w[i],
            row_h,
            cell,
            size=11,
            bold=bold,
            color=color,
        )
        x += col_w[i]
slide_footer(s, 12, TOTAL_SLIDES)


# ---- Slide 13 — Equipo ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)
slide_header(
    s,
    "12 · Equipo",
    "Quiénes lo construimos",
    "Experiencia técnica enterprise + conocimiento del mercado HR LATAM.",
)
team = [
    (
        "Ricardo Ascenda",
        "Founder · CEO · CTO",
        "Líder técnico y producto. Diseñó la arquitectura multi-tenant y la integración IA. Visión de negocio HR LATAM.",
    ),
    (
        "Backend Engineer",
        "NestJS · TypeORM · PostgreSQL",
        "Responsable de API, ciclos de evaluación, motor de objetivos y auditoría cross-tenant.",
    ),
    (
        "Full-stack Engineer",
        "Next.js · React · TypeScript",
        "Responsable de UI, PWA, dashboards de líder y RR.HH., experiencia mobile.",
    ),
]
tw = Inches(4.0)
th = Inches(3.0)
for i, (name, role, bio) in enumerate(team):
    left = Inches(0.6) + i * (tw + Inches(0.15))
    top = Inches(2.05)
    add_rect(s, left, top, tw, th, BG_PANEL, line_color=GOLD_DARK)
    # Avatar circle
    avatar = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.5), top + Inches(0.3), Inches(1.0), Inches(1.0))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = GOLD_PRIMARY
    avatar.line.color.rgb = GOLD_LIGHT
    initials_tf = avatar.text_frame
    initials_tf.margin_top = Emu(0)
    initials_tf.margin_bottom = Emu(0)
    initials_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = initials_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    initials = "".join([w[0] for w in name.split()[:2]]).upper()
    run = p.add_run()
    run.text = initials
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BG_BLACK
    add_text(
        s,
        left,
        top + Inches(1.4),
        tw,
        Inches(0.4),
        name,
        size=15,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left,
        top + Inches(1.8),
        tw,
        Inches(0.35),
        role,
        size=11,
        bold=True,
        color=GOLD_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        left + Inches(0.25),
        top + Inches(2.2),
        tw - Inches(0.4),
        Inches(0.8),
        bio,
        size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
add_rect(s, Inches(0.6), Inches(5.4), Inches(12.13), Inches(1.5), BG_PANEL, line_color=GOLD_DARK)
add_rect(s, Inches(0.6), Inches(5.4), Inches(0.06), Inches(1.5), GOLD_PRIMARY)
add_text(
    s,
    Inches(0.85),
    Inches(5.55),
    Inches(11.5),
    Inches(0.4),
    "POR QUÉ SOMOS LOS INDICADOS",
    size=12,
    bold=True,
    color=GOLD_LIGHT,
)
add_text(
    s,
    Inches(0.85),
    Inches(5.95),
    Inches(11.5),
    Inches(0.9),
    "Ya construimos lo difícil — multi-tenant, IA integrada y ciclos completos de evaluación.\nAhora aceleramos go-to-market regional con una plataforma diseñada para LATAM.",
    size=12,
    color=TEXT_MUTED,
    italic=True,
)
slide_footer(s, 13, TOTAL_SLIDES)


# ---- Slide 14 — Cierre ----
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, BG_BLACK)

# Decorative gold bar across top
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), GOLD_PRIMARY)

add_text(
    s,
    Inches(0.6),
    Inches(0.5),
    Inches(12),
    Inches(0.4),
    "13 · CIERRE",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_text(
    s,
    Inches(0.6),
    Inches(0.95),
    Inches(12),
    Inches(0.9),
    "EVA360 ya existe. Funciona. Lo difícil está hecho.",
    font=FONT_TITLE,
    size=32,
    bold=True,
    color=GOLD_LIGHT,
)
add_gold_bar(s, Inches(0.6), Inches(1.85), Inches(3.5))

# Two CTAs
add_rect(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(2.6), BG_PANEL, line_color=GOLD_DARK)
add_rect(s, Inches(0.6), Inches(2.2), Inches(0.08), Inches(2.6), GOLD_PRIMARY)
add_text(
    s,
    Inches(0.85),
    Inches(2.4),
    Inches(5.5),
    Inches(0.4),
    "PARA NUEVOS USUARIOS",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_text(
    s,
    Inches(0.85),
    Inches(2.85),
    Inches(5.5),
    Inches(1.8),
    "Empieza a evaluar a tu equipo de forma justa, rápida y con datos accionables. "
    "Demo gratuita en 15 minutos. Onboarding asistido en menos de una semana.",
    size=13,
    color=TEXT_MUTED,
)
add_text(
    s,
    Inches(0.85),
    Inches(4.3),
    Inches(5.5),
    Inches(0.4),
    "→ eva360.ascenda.cl/demo",
    size=13,
    bold=True,
    color=GOLD_LIGHT,
)

add_rect(s, Inches(6.85), Inches(2.2), Inches(6.0), Inches(2.6), BG_PANEL, line_color=GOLD_DARK)
add_rect(s, Inches(6.85), Inches(2.2), Inches(0.08), Inches(2.6), GOLD_PRIMARY)
add_text(
    s,
    Inches(7.1),
    Inches(2.4),
    Inches(5.5),
    Inches(0.4),
    "PARA INVERSIONISTAS",
    size=12,
    bold=True,
    color=GOLD_PRIMARY,
)
add_text(
    s,
    Inches(7.1),
    Inches(2.85),
    Inches(5.5),
    Inches(1.8),
    "Levantamos capital semilla para acelerar go-to-market en Chile + México "
    "y consolidar tracción regional en 18 meses. Producto listo, mercado validado.",
    size=13,
    color=TEXT_MUTED,
)
add_text(
    s,
    Inches(7.1),
    Inches(4.3),
    Inches(5.5),
    Inches(0.4),
    "→ ricardo@ascenda.cl",
    size=13,
    bold=True,
    color=GOLD_LIGHT,
)

# Contact strip
add_rect(s, Inches(0.6), Inches(5.2), Inches(12.13), Inches(1.5), GOLD_DARK)
add_text(
    s,
    Inches(0.6),
    Inches(5.45),
    Inches(12.13),
    Inches(0.5),
    "Build nos da el combustible para pasar de MVP demostrable a negocio sostenible.",
    size=14,
    italic=True,
    color=BG_BLACK,
    align=PP_ALIGN.CENTER,
)
add_text(
    s,
    Inches(0.6),
    Inches(6.0),
    Inches(12.13),
    Inches(0.5),
    "ricardo@ascenda.cl   ·   ascenda.cl   ·   eva360.ascenda.cl   ·   Santiago, Chile",
    size=14,
    bold=True,
    color=BG_BLACK,
    align=PP_ALIGN.CENTER,
)

slide_footer(s, 14, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUT_PPTX)
print(f"PPTX generado: {OUT_PPTX}")
print(f"Tamaño: {OUT_PPTX.stat().st_size / 1024:.1f} KB · Slides: {len(prs.slides)}")
